import os
import io
import sqlite3
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

# ----------------------
# App Config
# ----------------------
st.set_page_config(page_title="🏏 Cricsheet Analytics Suite", layout="wide")
st.title("🏏 Cricsheet Analytics Suite — EDA • Player/Team Insights • Exports")
DB_PATH = os.environ.get("CRICSHEET_DB", "scripts/cricket.db")

# ----------------------
# Utilities
# ----------------------
@st.cache_data(show_spinner=False)
def get_tables():
    with sqlite3.connect(DB_PATH) as conn:
        df = pd.read_sql("SELECT name FROM sqlite_master WHERE type='table'", conn)
    return sorted(df['name'].tolist())

@st.cache_data(show_spinner=False)
def read_sql(q: str) -> pd.DataFrame:
    with sqlite3.connect(DB_PATH) as conn:
        return pd.read_sql(q, conn)

@st.cache_data(show_spinner=False)
def load_matches(fmt: str) -> pd.DataFrame:
    tbl = f"{fmt}_matches"
    if tbl not in get_tables():
        return pd.DataFrame(columns=["match_id","format","teams","venue","date","toss_winner","match_winner"]) 
    df = read_sql(f"SELECT * FROM {tbl}")
    # normalize date/year
    if 'date' in df.columns:
        df['year'] = df['date'].astype(str).str[:4]
    return df

@st.cache_data(show_spinner=False)
def load_batting(fmt: str) -> pd.DataFrame:
    tbl = f"batting_stats_{fmt}"
    if tbl not in get_tables():
        return pd.DataFrame(columns=["batter","runs","ball","fours","sixes","strike_rate"]) 
    return read_sql(f"SELECT * FROM {tbl}")

# def load_batting(fmt: str) -> pd.DataFrame:
#     conn = sqlite3.connect("cricket.db")
#     query = f"""
#     SELECT 
#         batsman AS batter,
#         SUM(runs_batsman) AS runs,
#         COUNT(*) AS balls,
#         SUM(CASE WHEN runs_batsman = 4 THEN 1 ELSE 0 END) AS fours,
#         SUM(CASE WHEN runs_batsman = 6 THEN 1 ELSE 0 END) AS sixes,
#         ROUND(CAST(SUM(runs_batsman) AS FLOAT) / COUNT(*) * 100, 2) AS strike_rate
#     FROM {fmt}_matches
#     GROUP BY batsman
#     ORDER BY runs DESC;
#     """
#     df = pd.read_sql(query, conn)
#     conn.close()
#     return df

@st.cache_data(show_spinner=False)
def load_bowling(fmt: str) -> pd.DataFrame:
    tbl = f"bowling_stats_{fmt}"
    if tbl not in get_tables():
        return pd.DataFrame(columns=["bowler","ball","runs_conceded","wicket","overs","economy"]) 
    return read_sql(f"SELECT * FROM {tbl}")

@st.cache_data(show_spinner=False)
def load_team_results(fmt: str) -> pd.DataFrame:
    tbl = f"team_results_{fmt}"
    if tbl not in get_tables():
        return pd.DataFrame(columns=["team","wins"]) 
    return read_sql(f"SELECT * FROM {tbl}")

FORMATS = [fmt for fmt in ["odi","t20","test","ipl"] if f"{fmt}_matches" in get_tables()]
if not FORMATS:
    st.error("No match tables found in the database. Ensure cricket.db exists and tables are created.")
    st.stop()

# ----------------------
# Sidebar Filters
# ----------------------
with st.sidebar:
    st.header("⚙️ Controls")
    fmt = st.selectbox("Format", FORMATS, index=0, format_func=lambda s: s.upper())
    matches_df = load_matches(fmt)

    # Year filter from matches
    year_vals = sorted(matches_df['year'].dropna().unique().tolist()) if not matches_df.empty else []
    if year_vals:
        yr_min, yr_max = int(year_vals[0]), int(year_vals[-1])
        year_range = st.slider("Year Range", min_value=yr_min, max_value=yr_max, value=(yr_min, yr_max))
    else:
        year_range = None

    # Team filter from teams list string
    def expand_teams(s):
        try:
            # teams stored as string representation of list; eval safely
            s = s.strip()
            if s.startswith('[') and s.endswith(']'):
                out = s[1:-1].split(',')
                return [x.strip().strip("'").strip('"') for x in out]
        except Exception:
            pass
        return []

    teams_series = matches_df['teams'].dropna().apply(expand_teams)
    all_teams = sorted({t for lst in teams_series for t in lst}) if not matches_df.empty else []
    # team_filter = st.multiselect("Teams", options=all_teams, default=[])

    st.markdown("---")
    st.caption(f"Using database: **{DB_PATH}**")

# Apply filters to matches_df for EDA views
if not matches_df.empty and year_range:
    matches_df = matches_df[(matches_df['year'].astype(int) >= year_range[0]) & (matches_df['year'].astype(int) <= year_range[1])]
# if team_filter:
#     mask = matches_df['teams'].apply(lambda s: any(t in s for t in team_filter))
#     matches_df = matches_df[mask]

bat_df = load_batting(fmt)
bowl_df = load_bowling(fmt)
team_res_df = load_team_results(fmt)

# ----------------------
# Tabs
# ----------------------
t1, t2, t3, t4, t5, t6 = st.tabs([
    "📊 Overview",
    "🔍 EDA",
    "👤 Players",
    "🛡️ Teams",
    "📤 Exports",
    "🧭 Help"
])

# ----------------------
# Overview
# ----------------------
with t1:
    st.subheader(f"Overview — {fmt.upper()}")
    c1, c2, c3, c4 = st.columns(4)
    total_matches = len(matches_df)
    venues = matches_df['venue'].nunique() if 'venue' in matches_df.columns else 0
    teams_n = len(all_teams)
    winners_n = matches_df['match_winner'].nunique() if 'match_winner' in matches_df.columns else 0
    c1.metric("Matches", f"{total_matches:,}")
    c2.metric("Venues", f"{venues:,}")
    c3.metric("Teams", f"{teams_n:,}")
    c4.metric("Unique Winners", f"{winners_n:,}")

    # Matches per year
    if not matches_df.empty and 'year' in matches_df.columns:
        per_year = matches_df.groupby('year').size().reset_index(name='matches')
        fig = px.line(per_year, x='year', y='matches', markers=True, title=f"Matches per Year — {fmt.upper()}")
        st.plotly_chart(fig, use_container_width=True)

    # Top venues
    if not matches_df.empty:
        topv = matches_df['venue'].value_counts().head(10).reset_index()
        topv.columns = ['venue','matches']
        fig = px.bar(topv, x='matches', y='venue', orientation='h', title=f"Top 10 Venues — {fmt.upper()}")
        fig.update_layout(yaxis=dict(autorange='reversed'))
        st.plotly_chart(fig, use_container_width=True)

# ----------------------
# EDA
# ----------------------
with t2:
    st.subheader("Exploratory Data Analysis")
    if matches_df.empty:
        st.info("No matches to analyze with current filters.")
    else:
        c1, c2 = st.columns(2)
        # Toss vs Match winner
        tmp = matches_df.copy()
        tmp['outcome'] = np.where(tmp['toss_winner']==tmp['match_winner'], 'Toss = Match Winner', 'Toss ≠ Match Winner')
        # fig = px.bar(tmp['outcome'].value_counts().reset_index(), x='index', y='outcome',
        #              labels={'index':'Outcome','outcome':'Count'}, title='Toss vs Result')
        df_counts = tmp['outcome'].value_counts().reset_index()
        df_counts.columns = ['Outcome', 'Count']   # rename properly
        fig = px.bar(
            df_counts,
            x='Outcome',
            y='Count',
            labels={'Outcome': 'Match Outcome', 'Count': 'Count'},
            title='Toss vs Result',
            text='Count'
            )
        c1.plotly_chart(fig, use_container_width=True)

        # Winners (top 10)
        win = matches_df['match_winner'].value_counts().head(10).reset_index()
        win.columns = ['team','wins']
        fig2 = px.bar(win, x='wins', y='team', orientation='h', title='Top Winners (current filter)')
        fig2.update_layout(yaxis=dict(autorange='reversed'))
        c2.plotly_chart(fig2, use_container_width=True)

        # Heatmap Toss vs Match winner (seaborn)
        if matches_df['toss_winner'].notna().any() and matches_df['match_winner'].notna().any():
            ct = pd.crosstab(matches_df['toss_winner'], matches_df['match_winner'])
            fig_hm, ax = plt.subplots(figsize=(8,6))
            sns.heatmap(ct, ax=ax)
            ax.set_title("Heatmap: Toss Winner vs Match Winner")
            st.pyplot(fig_hm, use_container_width=True)

# ----------------------
# Players
# ----------------------
with t3:
    st.subheader(f"Player Performance — {fmt.upper()}")
    if bat_df.empty and bowl_df.empty:
        st.info("No player summary tables found. Build them with your fact-table script.")
    else:
        c1, c2 = st.columns(2)
        # Batting leaders
        if not bat_df.empty:
            min_balls = c1.slider("Min Balls (Batting)", 0, int(bat_df['ball'].max()) if 'ball' in bat_df.columns and not bat_df.empty else 1000, 300)
            bat_top = bat_df[bat_df['ball']>=min_balls].sort_values('runs', ascending=False).head(15)
            fig = px.bar(bat_top, x='runs', y='batter', orientation='h', title='Top Run-Scorers')
            fig.update_layout(yaxis=dict(autorange='reversed'))
            c1.plotly_chart(fig, use_container_width=True)

            # Strike rate leaders
            sr_top = bat_df[bat_df['ball']>=min_balls].sort_values('strike_rate', ascending=False).head(15)
            figsr = px.bar(sr_top, x='strike_rate', y='batter', orientation='h', title='Top Strike Rates')
            figsr.update_layout(yaxis=dict(autorange='reversed'))
            c1.plotly_chart(figsr, use_container_width=True)

        # Bowling leaders
        if not bowl_df.empty:
            min_balls_b = c2.slider("Min Balls (Bowling)", 0, int(bowl_df['ball'].max()) if 'ball' in bowl_df.columns and not bowl_df.empty else 1000, 300)
            bowl_top = bowl_df[bowl_df['ball']>=min_balls_b].sort_values('wicket', ascending=False).head(15)
            figb = px.bar(bowl_top, x='wicket', y='bowler', orientation='h', title='Top Wicket-Takers')
            figb.update_layout(yaxis=dict(autorange='reversed'))
            c2.plotly_chart(figb, use_container_width=True)

            # Wickets vs Economy
            scat = bowl_df[bowl_df['ball']>=min_balls_b]
            if not scat.empty:
                figsc = px.scatter(scat, x='wicket', y='economy', hover_name='bowler', title='Wickets vs Economy')
                c2.plotly_chart(figsc, use_container_width=True)

# ----------------------
# Teams
# ----------------------
with t4:
    st.subheader(f"Teams — {fmt.upper()}")
    if team_res_df.empty and matches_df.empty:
        st.info("No team tables available.")
    else:
        c1, c2 = st.columns(2)
        if not team_res_df.empty:
            # Top teams by wins
            topwins = team_res_df.sort_values('wins', ascending=False).head(15)
            figt = px.bar(topwins, x='wins', y='team', orientation='h', title='Team Wins (All-Time)')
            figt.update_layout(yaxis=dict(autorange='reversed'))
            c1.plotly_chart(figt, use_container_width=True)
        
        # Wins by year for selected teams
        if not matches_df.empty:
            sel_teams = c2.multiselect("Select teams for yearly trend", options=all_teams, default=all_teams[:2] if len(all_teams)>=2 else all_teams)
            if sel_teams:
                dfy = matches_df[matches_df['match_winner'].isin(sel_teams)].groupby(['year','match_winner']).size().reset_index(name='wins')
                figy = px.line(dfy, x='year', y='wins', color='match_winner', markers=True, title='Wins by Year (Selected Teams)')
                c2.plotly_chart(figy, use_container_width=True)

# ----------------------
# Exports (CSV + PPTX)
# ----------------------
with t5:
    st.subheader("Exports & Power BI")
    st.markdown("**CSV Exports for Power BI/Tableau**")
    c1, c2, c3 = st.columns(3)

    # Matches export
    if not matches_df.empty:
        c1.download_button(
            label="Download Matches CSV",
            data=matches_df.to_csv(index=False).encode('utf-8'),
            file_name=f"{fmt}_matches_filtered.csv",
            mime="text/csv"
        )
    if not bat_df.empty:
        c2.download_button(
            label="Download Batting CSV",
            data=bat_df.to_csv(index=False).encode('utf-8'),
            file_name=f"batting_stats_{fmt}.csv",
            mime="text/csv"
        )
    if not bowl_df.empty:
        c3.download_button(
            label="Download Bowling CSV",
            data=bowl_df.to_csv(index=False).encode('utf-8'),
            file_name=f"bowling_stats_{fmt}.csv",
            mime="text/csv"
        )

    st.markdown("---")
    st.markdown("**Generate PPTX from key charts**")

    def fig_to_png_bytes(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)
        return buf

    # Prepare matplotlib figures to embed in PPT
    figs = []
    # 1) Matches per year (mpl)
    if not matches_df.empty:
        per_year = matches_df.groupby('year').size().reset_index(name='matches')
        fig_mpl, ax = plt.subplots(figsize=(7,4))
        ax.plot(per_year['year'], per_year['matches'], marker='o')
        ax.set_title(f"Matches per Year — {fmt.upper()}")
        ax.tick_params(axis='x', rotation=45)
        figs.append((f"Matches per Year — {fmt.upper()}", fig_mpl))

    # 2) Top venues (mpl)
    if not matches_df.empty:
        topv = matches_df['venue'].value_counts().head(10).reset_index()
        topv.columns = ['venue','matches']
        fig_mpl2, ax2 = plt.subplots(figsize=(7,4))
        sns.barplot(data=topv, x='matches', y='venue', ax=ax2)
        ax2.set_title("Top Venues")
        figs.append(("Top Venues", fig_mpl2))

    # 3) Toss vs result (mpl)
    if not matches_df.empty:
        tmp = matches_df.copy()
        tmp['outcome'] = np.where(tmp['toss_winner']==tmp['match_winner'], 'Toss = Match Winner', 'Toss ≠ Match Winner')
        ct = tmp['outcome'].value_counts().reset_index()
        fig_mpl3, ax3 = plt.subplots(figsize=(6,4))
        ax3.bar(ct['outcome'], ct['count'])
        ax3.set_xlabel("Outcome")
        ax3.set_ylabel("Count")
        ax3.set_title("Toss vs Result")
        figs.append(("Toss vs Result", fig_mpl3))

    if st.button("Build PPTX"):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Cricsheet EDA Overview"
        slide.placeholders[1].text = f"Format: {fmt.upper()} | Years: {year_range if year_range else 'All'} | Teams: {', '.join(team_filter) if team_filter else 'All'}"

        for title, f in figs:
            img_stream = fig_to_png_bytes(f)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = Inches(1)
            top = Inches(1.2)
            slide.shapes.title.text = title
            slide.shapes.add_picture(img_stream, left, top, width=Inches(8))
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        st.download_button("Download Presentation (PPTX)", data=out, file_name=f"cricsheet_eda_{fmt}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    st.markdown("---")
    st.markdown("**Power BI** — Import CSVs above or connect directly to the SQLite DB via ODBC. Suggested tables: `*_matches`, `batting_stats_*`, `bowling_stats_*`, `team_results_*`. Create slicers for Year/Team/Format and build visuals (top scorers, wicket-takers, wins by team, toss vs result).")

# ----------------------
# Help
# ----------------------
with t6:
    st.subheader("Help & Notes")
    st.markdown(
        """
        **Data prerequisites**  
        • Database file: `cricket.db` in project root (change with env var `CRICSHEET_DB`).  
        • Tables expected: `{fmt}_matches`, `batting_stats_{fmt}`, `bowling_stats_{fmt}`, `team_results_{fmt}` for selected format.

        **Usage**  
        1) Choose **Format** in the sidebar, then restrict **Year Range** and **Teams** as needed.  
        2) Explore: Overview → EDA → Players → Teams.  
        3) Export: Download CSVs or click **Build PPTX** to generate a presentation of core charts.  

        **Tips**  
        • If tables are missing, run your ETL scripts to create fact tables (deliveries, batting, bowling, team results).  
        • For Power BI, either import the CSVs or set up a SQLite ODBC DSN and connect to `cricket.db` directly.  
        """
    )
